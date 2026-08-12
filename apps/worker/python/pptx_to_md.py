#!/usr/bin/env python3
"""确定性 PPTX 解析：python-pptx 逐页提取 → markdown（无模型、无网络）。
用法：python pptx_to_md.py <文件路径>  → markdown 打到 stdout。

每张幻灯片一节：标题占位符 → `## 标题`（没有标题的用 `## 幻灯片 N`），
其余文本框按形状顺序输出（保留项目符号层级为 markdown 列表缩进），表格转 markdown 表，
演讲者备注附在该页末尾（备注常含关键解释，对检索有价值）。

已知限制：
- 图片/图表/SmartArt 内的文字不抽取（走 kb-studio 的 vision 步骤）；
- 形状顺序按 shapes 集合顺序，不严格等于视觉阅读顺序（复杂排版可能次序略乱）;
- 组合形状(group)会递归展开。
"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from tabular_to_md import table_md  # 复用：丢空行/空列 + 跳标题横幅 + 单元格清洗


def _iter_shapes(shapes):
    """递归展开组合形状，按顺序产出叶子形状。"""
    for shp in shapes:
        if shp.shape_type == 6 and hasattr(shp, "shapes"):  # MSO_SHAPE_TYPE.GROUP
            for inner in _iter_shapes(shp.shapes):
                yield inner
        else:
            yield shp


def _text_frame_md(tf):
    """文本框 → markdown 行列表。段落缩进级别转成列表缩进；level 0 的独立成段。"""
    out = []
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            text = (para.text or "").strip()
        if not text:
            continue
        lvl = getattr(para, "level", 0) or 0
        if lvl > 0:
            out.append("  " * (lvl - 1) + "- " + text)
        else:
            out.append(text)
    return out


def pptx_to_md(path):
    from pptx import Presentation

    prs = Presentation(path)
    blocks = []
    for idx, slide in enumerate(prs.slides, start=1):
        # 标题：优先用标题占位符
        title = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except Exception:
            title = ""
        blocks.append(f"## {title}" if title else f"## 幻灯片 {idx}")

        # 用 shape_id（pptx 文件里的稳定标识）判重，不能用 id(obj)：
        # python-pptx 每次访问 shapes 都会新建包装对象，Python 对象身份对不上，
        # 结果就是标题被输出两遍。
        title_shape_id = None
        try:
            if slide.shapes.title is not None:
                title_shape_id = slide.shapes.title.shape_id
        except Exception:
            pass

        for shp in _iter_shapes(slide.shapes):
            if title_shape_id is not None and getattr(shp, "shape_id", None) == title_shape_id:
                continue  # 标题已单独输出，别重复
            if getattr(shp, "has_table", False):
                tbl = shp.table
                grid = [[c.text for c in row.cells] for row in tbl.rows]
                md = table_md(grid)
                if md:
                    blocks.append(md)
                continue
            if getattr(shp, "has_text_frame", False):
                lines = _text_frame_md(shp.text_frame)
                if lines:
                    blocks.append("\n".join(lines))

        # 演讲者备注
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    blocks.append(f"> 备注：{note}")
        except Exception:
            pass

    return "\n\n".join(b for b in blocks if b and b.strip())


def main():
    if len(sys.argv) < 2:
        sys.stderr.write("用法: python pptx_to_md.py <文件路径>\n")
        sys.exit(2)
    md = pptx_to_md(sys.argv[1])
    sys.stdout.write(md.strip() + "\n")


if __name__ == "__main__":
    main()
